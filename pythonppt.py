#!/usr/bin/env python
# coding: utf-8

# ### Criar PPT
# 
# python-pptx
# 
# Documentação: https://python-pptx.readthedocs.io/en/latest/index.html

# ### Criar PPT baseado em Modelo

# In[14]:


from pptx import Presentation

apresentacao = Presentation()

# editar o ppt 
slide1 = apresentacao.slides.add_slide(apresentacao.slide_layouts[0]) # slide com titulo e subtitulo
# slide1 = apresentacao.slides.add_slide(apresentacao.slide_layouts[6]) # slide em branco

titulo = slide1.shapes.title 
subtitulo = slide1.placeholders[1]

titulo.text = "1º Slide do Lira"
subtitulo.text = "Tamo criando ppt com Python"

# salvar esse ppt
apresentacao.save("MeuPPT.pptx")


# ### Criar PPT personalizados

# In[19]:


from pptx import Presentation
from pptx.util import Inches, Pt

apresentacao = Presentation()

# editar o ppt
slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])

x = Inches(1)
y = Inches(1)
largura = Inches(2)
altura = Inches(2)
caixa_texto = slide.shapes.add_textbox(x, y, largura, altura)

# 1 forma de editar uma caixa de texto
caixa_texto.text = "Vendas de Janeiro"

# 2 forma: criar parágrafos
text_frame = caixa_texto.text_frame
paragrafo = text_frame.add_paragraph()
paragrafo.text = "R$10.000"
paragrafo.font.bold = True
paragrafo.font.size = Pt(30)

paragrafo = text_frame.add_paragraph()
paragrafo.text = "Produto mais vendido: IPhone"

# salvar esse ppt
apresentacao.save("MeuPPT2.pptx")


# ### Criar PPTs com Gráficos

# In[20]:


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

apresentacao = Presentation()

slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])

# criar o nosso gráfico
produtos = ["IPhone", "IPad", "Airpod"] # categorias (eixo x)
vendas = [1500, 1000, 2000] # serie de dados (eixo y)

x = Inches(1)
y = Inches(1)
largura = Inches(5)
altura = Inches(3)

dados_grafico = CategoryChartData()
dados_grafico.categories = produtos
dados_grafico.add_series("Vendas", vendas)
# tipos de gráfico: https://python-pptx.readthedocs.io/en/latest/api/enum/XlChartType.html#xlcharttype
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, largura, altura, dados_grafico)

# salvar esse ppt
apresentacao.save("MeuPPT3.pptx")


# In[ ]:



